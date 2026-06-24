#!/usr/bin/env python3
"""Generator für die Workshop-Präsentation "KI im Unternehmen".

Erzeugt ein vollständig gestaltetes PowerPoint-Deck (.pptx) für einen Workshop
mit kleinen Unternehmern. Inhalt und Reihenfolge folgen dem roten Faden:
Vorstellung -> Digitale Transformation -> Prozessoptimierung -> Faktor Mensch
-> KI-gestütztes Coding -> Produkt-Showcase -> Angebot / Call to Action.

Design: helles, corporate-freundliches Theme (weißer Hintergrund, dunkler
Navy-Text, Teal/Orange als Akzente) — gut für helle Räume und als Handout.

Ausführen:
    pip install python-pptx
    python build_deck.py

Ergebnis: KI-Workshop.pptx im selben Verzeichnis.

Platzhalter wie [Ihr Name] und [Ihr Unternehmen] vor der Präsentation ersetzen.
Die genannten Kennzahlen sind branchenübliche Richtwerte (Quellen in den
Notizen) und sollten vor dem Vortrag kurz gegengeprüft / aktualisiert werden.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Design-System (helles / corporate Theme)
# ---------------------------------------------------------------------------
# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Farbpalette – hell, ruhig, hoher Kontrast für helle Räume & Druck
PAGE_BG   = RGBColor(0xFF, 0xFF, 0xFF)   # Seitenhintergrund weiß
SURFACE   = RGBColor(0xEE, 0xF2, 0xF7)   # Karten / Flächen (zartes Grau-Blau)
INK       = RGBColor(0x0B, 0x1B, 0x2B)   # Primärtext / dunkle Elemente (Navy)
SURFACE_DK = RGBColor(0x14, 0x26, 0x3B)  # dunkle Karte auf dunklem Slide
TEAL      = RGBColor(0x0D, 0x94, 0x88)   # Akzent 1 (kontraststark auf Weiß)
AMBER     = RGBColor(0xC2, 0x41, 0x0C)   # Akzent 2 (Burnt Orange)
MUTED     = RGBColor(0x51, 0x64, 0x7B)   # Sekundärtext auf Hell
MUTED_LT  = RGBColor(0xAE, 0xBE, 0xD0)   # Sekundärtext auf Dunkel
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Hilfsfunktionen
# ---------------------------------------------------------------------------
def _set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def _no_line(shape):
    shape.line.fill.background()


def add_rect(slide, x, y, w, h, color, line=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if not line:
        _no_line(shp)
    # Theme-Standardschatten abschalten -> flaches, sauberes Design.
    shp.shadow.inherit = False
    return shp


def add_round_rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    _no_line(shp)
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.05, space_after=Pt(6)):
    """runs: list of paragraphs; each paragraph is a list of (text, size, color, bold) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = space_after
        p.space_before = Pt(0)
        for (text, size, color, bold) in para:
            r = p.add_run()
            r.text = text
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
    return tb


def add_page_number(slide, n):
    add_text(slide, Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4),
             [[(str(n), 11, MUTED, False)]], align=PP_ALIGN.RIGHT)


def kicker(slide, text, color=TEAL, x=Inches(0.9), y=Inches(0.55)):
    add_text(slide, x, y, Inches(10), Inches(0.4),
             [[(text.upper(), 13, color, True)]])


_page = 0


def new_slide(bg=PAGE_BG, number=True):
    global _page
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, bg)
    _page += 1
    if number:
        add_page_number(s, _page)
    return s


# ---------------------------------------------------------------------------
# Slide-Vorlagen
# ---------------------------------------------------------------------------
def title_slide():
    global _page
    s = prs.slides.add_slide(BLANK)
    _set_bg(s, PAGE_BG)
    _page += 1
    # Akzentbalken links
    add_rect(s, 0, 0, Inches(0.35), SLIDE_H, TEAL)
    # dezenter Block rechts oben
    add_rect(s, Inches(11.3), 0, Inches(2.03), Inches(0.35), AMBER)

    add_text(s, Inches(0.95), Inches(1.6), Inches(11.5), Inches(0.5),
             [[("WORKSHOP · KI FÜR KLEINE UNTERNEHMEN", 15, TEAL, True)]])
    add_text(s, Inches(0.95), Inches(2.25), Inches(11.6), Inches(2.6),
             [[("KI im Unternehmen:", 54, INK, True)],
              [("von der Idee zur Umsetzung", 54, INK, True)]],
             line_spacing=1.0, space_after=Pt(2))
    add_text(s, Inches(0.98), Inches(4.5), Inches(11.4), Inches(0.8),
             [[("Wie KI und KI-gestütztes Coding Ihrem Betrieb Zeit, "
                "Geld und Nerven sparen.", 20, MUTED, False)]])
    # Presenter
    add_rect(s, Inches(0.98), Inches(5.7), Inches(0.5), Inches(0.06), AMBER)
    add_text(s, Inches(0.98), Inches(5.9), Inches(11), Inches(0.9),
             [[("[Ihr Name]", 20, INK, True)],
              [("[Ihr Unternehmen] · KI- & Digitalisierungsberatung", 15, MUTED, False)]])


def section_slide(no, title, subtitle):
    s = new_slide(PAGE_BG)
    # dunkle Bandfläche für visuellen Rhythmus
    add_rect(s, 0, Inches(2.55), SLIDE_W, Inches(2.4), INK)
    add_text(s, Inches(0.95), Inches(2.05), Inches(3), Inches(1.2),
             [[(no, 110, TEAL, True)]])
    add_text(s, Inches(3.1), Inches(2.95), Inches(9.4), Inches(1.0),
             [[(title, 40, WHITE, True)]], anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(3.13), Inches(4.0), Inches(9.0), Inches(1.2),
             [[(subtitle, 18, MUTED_LT, False)]])
    return s


def content_slide(kick, heading):
    s = new_slide(PAGE_BG)
    kicker(s, kick)
    add_text(s, Inches(0.9), Inches(0.95), Inches(11.6), Inches(1.0),
             [[(heading, 32, INK, True)]])
    add_rect(s, Inches(0.95), Inches(1.85), Inches(0.9), Inches(0.06), TEAL)
    return s


def bullets(slide, items, x=Inches(0.95), y=Inches(2.2), w=Inches(11.4),
            h=Inches(4.6), size=19, gap=Pt(14), lead_color=INK):
    """items: list of (lead, rest) – lead fett, rest gedämpft."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = gap
        p.line_spacing = 1.1
        # Bullet-Marker
        r0 = p.add_run()
        r0.text = "▸  "
        r0.font.name = FONT
        r0.font.size = Pt(size)
        r0.font.color.rgb = TEAL
        r0.font.bold = True
        lead, rest = item
        if lead:
            r1 = p.add_run()
            r1.text = lead
            r1.font.name = FONT
            r1.font.size = Pt(size)
            r1.font.color.rgb = lead_color
            r1.font.bold = True
        if rest:
            r2 = p.add_run()
            r2.text = (" — " if lead else "") + rest
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = MUTED
            r2.font.bold = False
    return tb


def card(slide, x, y, w, h, title, body, accent=TEAL):
    add_round_rect(slide, x, y, w, h, SURFACE)
    add_rect(slide, x, y, Inches(0.09), h, accent)
    add_text(slide, x + Inches(0.35), y + Inches(0.25), w - Inches(0.6),
             Inches(0.6), [[(title, 17, INK, True)]])
    add_text(slide, x + Inches(0.35), y + Inches(0.85), w - Inches(0.6),
             h - Inches(1.0), [[(body, 13.5, MUTED, False)]], line_spacing=1.08)


def stat(slide, x, y, w, big, label, accent=TEAL):
    add_text(slide, x, y, w, Inches(1.1),
             [[(big, 60, accent, True)]], align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.15), w, Inches(1.1),
             [[(label, 15, INK, False)]], align=PP_ALIGN.CENTER,
             line_spacing=1.05)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ===========================================================================
# 1 · TITEL
# ===========================================================================
title_slide()

# ===========================================================================
# 2 · AGENDA
# ===========================================================================
s = content_slide("Überblick", "Unser Weg durch den Workshop")
agenda = [
    ("1 · Wer ich bin", "Hintergrund und was ich bereits gebaut habe"),
    ("2 · Digitale Transformation", "Was KI heute leistet — in Zahlen"),
    ("3 · Prozessoptimierung", "Erst verstehen, dann automatisieren"),
    ("4 · Der Faktor Mensch", "Warum Technik der einfache Teil ist"),
    ("5 · KI-gestütztes Coding", "Live-Demo: Software in Minuten statt Monaten"),
    ("6 · Was als Nächstes", "Wie Sie das in Ihrem Betrieb nutzen"),
]
bullets(s, agenda, y=Inches(2.15), gap=Pt(13), size=19)
notes(s, "Kurz durch die Agenda führen (max. 1 Minute). Versprechen: Am Ende "
          "wissen Sie, WO KI in Ihrem Betrieb hilft und SEHEN, wie schnell "
          "Software heute entsteht. Ziel ist nicht, dass Sie selbst Entwickler "
          "werden — sondern dass Sie die Chancen erkennen.")

# ===========================================================================
# 3 · VORSTELLUNG
# ===========================================================================
s = content_slide("Vorstellung", "Wer vor Ihnen steht")
bullets(s, [
    ("[Ihr Name]", "Ich entwickle und berate an der Schnittstelle von KI, "
     "Software und Geschäftsprozessen."),
    ("Macher, nicht nur Berater", "Ich baue selbst Produkte — ich rede nicht "
     "nur über Technologie, ich setze sie ein."),
    ("Fokus", "Digitale Transformation: echte Probleme im Betrieb finden und "
     "mit dem richtigen Werkzeug lösen."),
    ("Heute", "Ich zeige Ihnen, was möglich ist — und wo der Hebel für Ihr "
     "Unternehmen liegt."),
], y=Inches(2.2), gap=Pt(18), size=20)
notes(s, "Persönlich werden. Eine kurze Geschichte, warum Sie das machen. "
          "Betonen: Sie sind Praktiker. Glaubwürdigkeit entsteht durch die "
          "Produkte auf der nächsten Folie.")

# ===========================================================================
# 4 · PRODUKTE (Showcase früh, als Beleg)
# ===========================================================================
s = content_slide("Referenzen", "Was ich bereits gebaut habe")
cw, ch = Inches(3.78), Inches(1.55)
gx, gy = Inches(0.4), Inches(0.32)
x0, y0 = Inches(0.95), Inches(2.15)
items = [
    ("Video-Meeting-Plattform", "Eine Art „Google Meet“ inkl. Meeting-Bot, der "
     "Gespräche begleitet.", TEAL),
    ("Startup-Ökosystem", "Plattform, die Gründer, Wissen und Ressourcen "
     "vernetzt.", AMBER),
    ("„Learn how to talk“", "Trainingsapp für Kommunikation und sicheres "
     "Auftreten.", TEAL),
    ("Job-Interview-Training", "KI-gestützte Vorbereitung auf "
     "Bewerbungsgespräche.", AMBER),
    ("Eigenes Game", "Komplettes Spiel — selbst entwickelt.", TEAL),
    ("… und weitere Produkte", "Vom Prototyp bis zum laufenden System.", AMBER),
]
for i, (t, b, acc) in enumerate(items):
    col = i % 3
    row = i // 3
    x = x0 + col * (cw + gx)
    y = y0 + row * (ch + gy)
    card(s, x, y, cw, ch, t, b, acc)
notes(s, "Das ist Ihr Vertrauens-Anker. Kurz pro Produkt 1 Satz. Botschaft: "
          "„All das wurde mit modernen KI-Werkzeugen gebaut — von einer "
          "kleinen Mannschaft, schnell.“ Damit ist die Brücke zum Rest des "
          "Workshops geschlagen.")

# ===========================================================================
# 5 · SECTION: DIGITALE TRANSFORMATION
# ===========================================================================
section_slide("01", "Digitale Transformation",
              "Was KI heute wirklich leistet — und warum jetzt der richtige "
              "Zeitpunkt ist.")

# 6 · Was bedeutet das
s = content_slide("Grundlagen", "Digitale Transformation ist kein IT-Projekt")
bullets(s, [
    ("Es geht um Wertschöpfung", "nicht um Technik um der Technik willen. "
     "Werkzeuge sind Mittel zum Zweck."),
    ("Repetitive Arbeit", "verschwindet — Angebote, Berichte, Datenpflege, "
     "Standard-Kommunikation."),
    ("Wissen wird zugänglich", "Antworten in Sekunden statt Suchen in Ordnern."),
    ("Kleine Betriebe im Vorteil", "Sie sind schneller und können neue "
     "Werkzeuge sofort einsetzen — ohne Konzern-Bürokratie."),
], y=Inches(2.2), gap=Pt(16), size=19)
notes(s, "Klarstellen: Digitalisierung ≠ neue Software kaufen. Es ist eine "
          "Frage, WIE gearbeitet wird. Kleine Unternehmen sind agil — das ist "
          "ihr Trumpf gegenüber großen Wettbewerbern.")

# 7 · KI in Zahlen
s = content_slide("Die Zahlen", "Warum sich das Hinschauen lohnt")
add_round_rect(s, Inches(0.95), Inches(2.25), Inches(11.45), Inches(2.4), SURFACE)
stat(s, Inches(1.0), Inches(2.6), Inches(3.7), "bis 40 %", "weniger Zeit für "
     "Routineaufgaben durch KI-Assistenz", TEAL)
stat(s, Inches(4.85), Inches(2.6), Inches(3.7), "10×", "schnellere Erstellung "
     "von Texten, Angeboten & Code", AMBER)
stat(s, Inches(8.7), Inches(2.6), Inches(3.7), "24/7", "verfügbar — KI macht "
     "keine Pause und wird nicht krank", TEAL)
add_text(s, Inches(0.95), Inches(5.05), Inches(11.45), Inches(1.6),
         [[("Wichtig: Der Gewinn entsteht nicht durch „KI haben“, sondern "
            "durch KI an der richtigen Stelle. Ein gut gewählter Anwendungsfall "
            "amortisiert sich oft in wenigen Wochen.", 16, MUTED, False)]])
notes(s, "Zahlen sind branchenübliche Richtwerte (u. a. McKinsey/„The economic "
          "potential of generative AI“, Studien zu Entwickler-Produktivität mit "
          "KI-Assistenten). VOR DEM VORTRAG kurz aktualisieren und als „grobe "
          "Größenordnung“ präsentieren — nicht als exakte Garantie. Die "
          "eigentliche Botschaft steht unten: richtige Stelle schlägt "
          "Technik-Hype.")

# 8 · Warum jetzt
s = content_slide("Timing", "Warum gerade jetzt der richtige Moment ist")
bullets(s, [
    ("Die Werkzeuge sind reif", "Was vor zwei Jahren Forschung war, ist heute "
     "Alltag und bezahlbar."),
    ("Die Einstiegshürde ist gefallen", "Sie brauchen kein eigenes "
     "IT-Team mehr für vieles."),
    ("Der Wettbewerb beginnt gerade", "Wer jetzt lernt, hat einen echten "
     "Vorsprung — in 2 Jahren ist es Standard."),
    ("Risiko des Abwartens", "Nicht „zu früh investieren“, sondern „zu spät "
     "anfangen“ ist die eigentliche Gefahr."),
], y=Inches(2.2), gap=Pt(16), size=19)
notes(s, "Dringlichkeit erzeugen, aber ohne Angstmache. Framing: Frühe, kleine "
          "Schritte schlagen großes Abwarten. Überleitung zu Abschnitt 2: "
          "„Bevor wir Werkzeuge einsetzen, müssen wir wissen WO.“")

# ===========================================================================
# 9 · SECTION: PROZESSOPTIMIERUNG
# ===========================================================================
section_slide("02", "Prozesse verstehen",
              "Der wichtigste Schritt kommt vor der Technik: zu wissen, wo es "
              "im Betrieb wirklich klemmt.")

# 10 · Erst verstehen
s = content_slide("Reihenfolge", "Erst verstehen — dann automatisieren")
bullets(s, [
    ("Häufigster Fehler", "Mit dem Werkzeug starten, bevor das Problem klar "
     "ist. Dann digitalisiert man Chaos."),
    ("Richtige Reihenfolge", "Prozess verstehen → Engpass finden → richtige "
     "Lösung wählen → erst dann umsetzen."),
    ("„Automatisieren Sie keinen schlechten Prozess“", "sonst machen Sie den "
     "Fehler nur schneller."),
    ("Genau hier komme ich ins Spiel", "Den Betrieb verstehen, bevor "
     "irgendjemand eine Zeile Code schreibt."),
], y=Inches(2.2), gap=Pt(16), size=19)
notes(s, "Kernbotschaft des ganzen Workshops. Beispiel erzählen: Ein Betrieb "
          "wollte „eine App“, brauchte aber eigentlich nur eine saubere "
          "Vorlage + Automatisierung. Das ist Ihr Mehrwert als Berater.")

# 11 · Potenziale finden
s = content_slide("Analyse", "Wo liegen Ihre Potenziale?")
qa = [
    ("Was wiederholt sich täglich?", "Jede manuelle Routine ist ein Kandidat."),
    ("Wo entstehen Wartezeiten?", "Übergaben, Freigaben, Suchen nach Infos."),
    ("Wo passieren Fehler?", "Doppelte Eingaben, Copy-Paste, Medienbrüche."),
    ("Was raubt Ihnen Energie?", "Aufgaben, die niemand gern macht."),
    ("Wo sagen Kunden „zu langsam“?", "Reaktionszeit ist oft der Engpass."),
    ("Was hängt an einer Person?", "Wissen im Kopf ist ein Risiko."),
]
cw, ch = Inches(5.6), Inches(1.18)
for i, (t, b) in enumerate(qa):
    col = i % 2
    row = i // 2
    x = Inches(0.95) + col * (cw + Inches(0.25))
    y = Inches(2.15) + row * (ch + Inches(0.2))
    card(s, x, y, cw, ch, t, b, AMBER if col else TEAL)
notes(s, "Das sind Ihre Diagnose-Fragen. Publikum aktivieren: „Denken Sie an "
          "EINEN Prozess in Ihrem Betrieb, der Sie nervt.“ Diese Folie ist "
          "praktisch und einprägsam — evtl. als Handout mitgeben.")

# 12 · Vom Problem zur Lösung
s = content_slide("Vorgehen", "Vom Problem zur Lösung — in 4 Schritten")
steps = [
    ("1", "Beobachten", "Prozess Schritt für Schritt aufnehmen — so wie er "
     "wirklich läuft.", TEAL),
    ("2", "Bewerten", "Wo ist der größte Hebel? Aufwand vs. Wirkung.", AMBER),
    ("3", "Lösen", "Richtiges Werkzeug wählen — oft KI, manchmal nur eine "
     "Vorlage.", TEAL),
    ("4", "Einführen", "Klein starten, messen, Menschen mitnehmen.", AMBER),
]
cw = Inches(2.78)
for i, (n, t, b, acc) in enumerate(steps):
    x = Inches(0.95) + i * (cw + Inches(0.18))
    add_round_rect(s, x, Inches(2.3), cw, Inches(3.0), SURFACE)
    add_text(s, x, Inches(2.55), cw, Inches(1.0), [[(n, 48, acc, True)]],
             align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.2), Inches(3.55), cw - Inches(0.4), Inches(0.5),
             [[(t, 18, INK, True)]], align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), Inches(4.05), cw - Inches(0.5), Inches(1.2),
             [[(b, 13, MUTED, False)]], align=PP_ALIGN.CENTER, line_spacing=1.08)
notes(s, "Einfache, wiederholbare Methode. Betonen: Schritt 3 (Werkzeug) ist "
          "der KLEINSTE Teil. Die Arbeit steckt in 1, 2 und 4. Überleitung zu "
          "Abschnitt 3: Schritt 4 — die Menschen — ist der schwerste.")

# ===========================================================================
# 13 · SECTION: FAKTOR MENSCH
# ===========================================================================
section_slide("03", "Der Faktor Mensch",
              "Technik ist der einfache Teil. Menschen für Veränderung zu "
              "gewinnen, entscheidet über Erfolg oder Scheitern.")

# 14 · Technik vs Mensch
s = content_slide("Die eigentliche Hürde", "Technik ist einfach — Menschen sind schwer")
add_round_rect(s, Inches(0.95), Inches(2.2), Inches(5.55), Inches(3.6), SURFACE)
add_rect(s, Inches(0.95), Inches(2.2), Inches(5.55), Inches(0.09), TEAL)
add_text(s, Inches(1.25), Inches(2.45), Inches(5.0), Inches(0.6),
         [[("Technik", 22, TEAL, True)]])
bullets(s, [
    ("Verfügbar", "Werkzeuge gibt es im Überfluss."),
    ("Schnell", "Einrichtung dauert oft Minuten."),
    ("Planbar", "Funktioniert oder funktioniert nicht."),
], x=Inches(1.25), y=Inches(3.15), w=Inches(5.0), size=16, gap=Pt(12))

add_round_rect(s, Inches(6.85), Inches(2.2), Inches(5.55), Inches(3.6), SURFACE)
add_rect(s, Inches(6.85), Inches(2.2), Inches(5.55), Inches(0.09), AMBER)
add_text(s, Inches(7.15), Inches(2.45), Inches(5.0), Inches(0.6),
         [[("Menschen", 22, AMBER, True)]])
bullets(s, [
    ("Gewohnheiten", "„Das haben wir immer so gemacht.“"),
    ("Angst", "„Ersetzt mich die KI?“ — muss ernst genommen werden."),
    ("Vertrauen", "entsteht nur durch Erfolge, die man selbst erlebt."),
], x=Inches(7.15), y=Inches(3.15), w=Inches(5.0), size=16, gap=Pt(12))
notes(s, "Die wichtigste Folie für Ihr Verkaufsargument. Tools kann jeder "
          "googeln — den Wandel beim TEAM zu schaffen, ist die eigentliche "
          "Leistung. Genau dafür braucht man jemanden wie Sie.")

# 15 · Akzeptanz schaffen
s = content_slide("Change", "So gewinnen Sie Ihr Team")
bullets(s, [
    ("Nutzen zuerst zeigen", "Menschen ändern sich, wenn das Werkzeug IHR "
     "Leben leichter macht — nicht das der Geschäftsführung."),
    ("Klein anfangen", "Ein Prozess, ein sichtbarer schneller Erfolg, dann "
     "ausweiten."),
    ("Mitmachen lassen", "Wer mitgestaltet, blockiert nicht. Die "
     "Erfahrensten zu Vorbildern machen."),
    ("Ehrlich kommunizieren", "KI nimmt Aufgaben ab, nicht Menschen weg — "
     "Zeit für das Wesentliche."),
    ("Begleiten statt anordnen", "Einführung braucht Geduld und jemanden, der "
     "dranbleibt."),
], y=Inches(2.1), gap=Pt(13), size=18)
notes(s, "Praktischer Werkzeugkasten für Akzeptanz. Pointe: „Die beste "
          "Software ist wertlos, wenn sie keiner benutzt.“ Überleitung: „Jetzt "
          "der spannende Teil — wie Software heute überhaupt entsteht.“")

# ===========================================================================
# 16 · SECTION: KI-CODING
# ===========================================================================
section_slide("04", "KI-gestütztes Coding",
              "Wie aus einer Idee in Minuten funktionierende Software wird — "
              "live.")

# 17 · Was ist möglich
s = content_slide("Was heute geht", "Software bauen, ohne Entwickler zu sein")
bullets(s, [
    ("Sie beschreiben — die KI baut", "In normaler Sprache sagen, was Sie "
     "wollen. Die KI schreibt den Code."),
    ("Von Stunden, nicht Monaten", "Prototypen entstehen am selben Tag, nicht "
     "im nächsten Quartal."),
    ("Bruchteil der Kosten", "Was früher ein Team und ein Budget brauchte, "
     "schafft heute eine Person."),
    ("Genau so sind meine Produkte entstanden", "Die Beispiele vom Anfang — "
     "alle mit diesen Methoden gebaut."),
], y=Inches(2.2), gap=Pt(16), size=19)
notes(s, "Erwartung setzen für die Demo. WICHTIG: Ehrlich bleiben — KI ersetzt "
          "keinen erfahrenen Kopf bei komplexen / kritischen Systemen. Aber für "
          "80 % der Alltagsfälle ist sie ein Gamechanger. Das verkauft Ihre "
          "Rolle: Sie sind der erfahrene Kopf.")

# 18 · Setup
s = content_slide("Das Setup", "Womit ich gleich arbeite")
cards3 = [
    ("Der KI-Assistent", "Ein Coding-Agent, der in normaler Sprache Anweisungen "
     "versteht und direkt umsetzt.", TEAL),
    ("Der Arbeitsplatz", "Editor + Vorschau: links die Anweisung, rechts das "
     "Ergebnis in Echtzeit.", AMBER),
    ("Der Kreislauf", "Sagen → sehen → verbessern. In kurzen Schleifen zum "
     "fertigen Ergebnis.", TEAL),
]
cw = Inches(3.78)
for i, (t, b, acc) in enumerate(cards3):
    x = Inches(0.95) + i * (cw + Inches(0.18))
    card(s, x, Inches(2.4), cw, Inches(3.0), t, b, acc)
notes(s, "Kurz halten — nicht in Technik abdriften. Das Publikum muss den "
          "Ablauf verstehen, nicht die Werkzeuge im Detail. Dann auf die "
          "nächste Folie: LIVE.")

# 19 · DEMO BREAK
s = new_slide(TEAL, number=False)
add_text(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.4),
         [[("LIVE-DEMO", 24, WHITE, True)]])
add_text(s, Inches(1.0), Inches(3.1), Inches(11.3), Inches(1.6),
         [[("Wir bauen jetzt gemeinsam etwas.", 44, WHITE, True)]])
add_text(s, Inches(1.02), Inches(4.6), Inches(11.0), Inches(1.0),
         [[("Schauen Sie auf den Ablauf, nicht auf den Code: "
            "beschreiben → erzeugen → verbessern.", 18, WHITE, False)]])
add_rect(s, Inches(1.02), Inches(2.35), Inches(1.2), Inches(0.08), INK)
notes(s, "AUF SCREEN-SHARE / EDITOR WECHSELN. Demo-Tipps: 1) Etwas wählen, das "
          "die Zielgruppe versteht (z. B. einfache Kundenanfrage-Seite oder "
          "Angebots-Generator). 2) Laut denken. 3) Bewusst auch eine kleine "
          "Korrektur machen — zeigt, dass man im Dialog steuert. 4) BACKUP "
          "bereithalten (Screenshots/Aufnahme), falls Live nicht klappt oder "
          "WLAN schwächelt. Zeitrahmen: 5–8 Minuten.")

# 20 · Demo recap
s = content_slide("Einordnung", "Was Sie gerade gesehen haben")
bullets(s, [
    ("Keine Zeile selbst getippt", "Idee in Sprache → fertiges Ergebnis."),
    ("Minuten statt Wochen", "Diese Geschwindigkeit verändert, was sich für "
     "einen kleinen Betrieb überhaupt lohnt."),
    ("Aber: jemand muss steuern", "Wissen, WAS man baut und WARUM, bleibt "
     "entscheidend — genau das ist die eigentliche Arbeit."),
    ("Werkzeug ≠ Lösung", "Die KI baut, was Sie sagen. Was sinnvoll ist, "
     "entscheidet das Verständnis für Ihren Betrieb."),
], y=Inches(2.2), gap=Pt(16), size=19)
notes(s, "Den „Wow“-Moment in Wert ummünzen. Roter Faden schließen: Technik ist "
          "schnell und billig geworden — der Engpass ist VERSTEHEN (Prozesse) "
          "und MENSCHEN. Beides liefern Sie.")

# ===========================================================================
# 21 · SECTION: NÄCHSTE SCHRITTE
# ===========================================================================
section_slide("05", "Ihr nächster Schritt",
              "Von der Inspiration zur Umsetzung — ohne dass Sie selbst zum "
              "Entwickler werden müssen.")

# 22 · Wie ich helfe
s = content_slide("Zusammenarbeit", "Wie ich Sie unterstütze")
cards4 = [
    ("Potenzial-Analyse", "Gemeinsam Ihre Prozesse durchgehen und die größten "
     "Hebel finden.", TEAL),
    ("Umsetzung", "Ich baue die Lösung — schnell, mit modernen KI-Methoden.",
     AMBER),
    ("Team-Einführung", "Ihre Mitarbeiter mitnehmen, damit die Lösung wirklich "
     "genutzt wird.", TEAL),
    ("Begleitung", "Dranbleiben, messen, nachschärfen — bis es im Alltag "
     "läuft.", AMBER),
]
cw, ch = Inches(5.6), Inches(1.45)
for i, (t, b, acc) in enumerate(cards4):
    col = i % 2
    row = i // 2
    x = Inches(0.95) + col * (cw + Inches(0.25))
    y = Inches(2.2) + row * (ch + Inches(0.25))
    card(s, x, y, cw, ch, t, b, acc)
notes(s, "Ihr Angebot, klar strukturiert. Botschaft an die Unternehmer: „Sie "
          "müssen das nicht selbst lernen. Sie brauchen jemanden, der versteht "
          "UND umsetzt.“ Das ist genau die Buchung, auf die der Workshop "
          "hinarbeitet.")

# 23 · Call to action (dunkler Abschluss für visuellen Punkt)
s = new_slide(INK, number=False)
add_rect(s, 0, 0, Inches(0.35), SLIDE_H, TEAL)
add_text(s, Inches(0.95), Inches(1.5), Inches(11.4), Inches(0.5),
         [[("LASSEN SIE UNS REDEN", 15, TEAL, True)]])
add_text(s, Inches(0.95), Inches(2.1), Inches(11.4), Inches(1.6),
         [[("Wo klemmt es bei Ihnen?", 46, WHITE, True)]])
add_text(s, Inches(0.98), Inches(3.5), Inches(11.0), Inches(1.0),
         [[("15 Minuten, ein Prozess, eine ehrliche Einschätzung — "
            "kostenlos und unverbindlich.", 19, MUTED_LT, False)]])
add_round_rect(s, Inches(0.98), Inches(4.7), Inches(5.4), Inches(1.6), SURFACE_DK)
add_text(s, Inches(1.3), Inches(4.95), Inches(5.0), Inches(1.2),
         [[("[Ihr Name]", 20, WHITE, True)],
          [("[E-Mail-Adresse]", 15, MUTED_LT, False)],
          [("[Telefon / Website]", 15, MUTED_LT, False)]], space_after=Pt(4))
add_round_rect(s, Inches(6.6), Inches(4.7), Inches(5.8), Inches(1.6), TEAL)
add_text(s, Inches(6.9), Inches(5.05), Inches(5.2), Inches(1.0),
         [[("Termin sichern", 22, WHITE, True)],
          [("Scannen Sie den QR-Code / schreiben Sie mir heute.", 14,
            WHITE, False)]], space_after=Pt(6))
notes(s, "Klarer, niedrigschwelliger Call to Action. Tipp: Hier einen echten "
          "QR-Code zu Ihrem Kalender (z. B. Calendly) einfügen und die "
          "Platzhalter mit echten Kontaktdaten ersetzen. Letzter Satz im "
          "Vortrag: Einladung zum kostenlosen Erstgespräch — nicht verkaufen, "
          "sondern helfen anbieten.")

# ---------------------------------------------------------------------------
out = "KI-Workshop.pptx"
prs.save(out)
print(f"Gespeichert: {out}  ({len(prs.slides)} Folien)")
